"""
EDIP - Executive Report Generator
=====================================
Takes the AI Consultant's structured JSON answers and renders a PowerPoint-ready
executive summary deck -- the "write executive summaries / produce PowerPoint-ready
reports" capability of the autonomous agent.

Uses python-pptx (simple, dependency-light, ideal for automated/scheduled report
generation in a production pipeline, as opposed to a one-off hand-crafted deck).

Run:
    python3 reports/generate_executive_report.py
Output:
    outputs/EDIP_Executive_Report.pptx
"""
import sys
from io import BytesIO
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))
OUT_DIR = ROOT / "outputs"
OUT_DIR.mkdir(exist_ok=True)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ModuleNotFoundError:
    Presentation = None
    Inches = Pt = Emu = RGBColor = PP_ALIGN = None
    PPTX_AVAILABLE = False

from agent.ai_consultant import AIConsultant

if PPTX_AVAILABLE:
    NAVY = RGBColor(0x10, 0x1B, 0x3A)
    ACCENT = RGBColor(0x2E, 0x6B, 0xE0)
    GREY = RGBColor(0x55, 0x5B, 0x66)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFB)
else:
    NAVY = ACCENT = GREY = WHITE = LIGHT_BG = None


def add_title_slide(prs, title, subtitle):
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required to build PPT reports. Install it with 'pip install python-pptx'.")

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(1.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = subtitle
    p2.font.size = Pt(18); p2.font.color.rgb = RGBColor(0xC7, 0xD2, 0xF0)
    return slide


def add_section_header(slide, text, top=None):
    if top is None:
        top = Inches(0.4)
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    return tb


def add_body_slide(prs, header):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, header)
    return slide


def add_bullets(slide, items, left=None, top=None, width=None, height=None, size=16):
    if left is None:
        left = Inches(0.6)
    if top is None:
        top = Inches(1.3)
    if width is None:
        width = Inches(12)
    if height is None:
        height = Inches(5.5)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = GREY
        p.space_after = Pt(10)
    return tb


def add_kpi_cards(slide, kpis, top=None):
    if top is None:
        top = Inches(1.3)
    n = len(kpis)
    card_w = Inches(2.7)
    gap = Inches(0.3)
    total_w = card_w * n + gap * (n - 1)
    x = Inches(0.6)
    for label, value in kpis:
        card = slide.shapes.add_shape(1, x, top, card_w, Inches(1.4))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xDD, 0xE2, 0xEE); card.line.width = Pt(1)
        card.shadow.inherit = False
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(10); tf.margin_top = Pt(10)
        p1 = tf.paragraphs[0]; p1.text = str(value)
        p1.font.size = Pt(24); p1.font.bold = True; p1.font.color.rgb = ACCENT
        p2 = tf.add_paragraph(); p2.text = label
        p2.font.size = Pt(12); p2.font.color.rgb = GREY
        x = Emu(int(x) + int(card_w) + int(gap))


def add_answer_summary_slide(prs, answer):
    slide = add_body_slide(prs, "Executive Summary")
    items = []
    if answer.get("question"):
        items.append(f"Question: {answer['question']}")
    if answer.get("narrative"):
        items.append(f"Summary: {answer['narrative']}")
    if answer.get("quantitative_findings"):
        items.append("Quantitative findings:")
        for key, value in answer["quantitative_findings"].items():
            items.append(f"  • {key.replace('_', ' ').title()}: {value}")
    if answer.get("summary"):
        items.append("Highlights:")
        for key, value in answer["summary"].items():
            items.append(f"  • {key.replace('_', ' ').title()}: {value}")
    if answer.get("recommended_actions"):
        items.append("Recommended actions:")
        for action in answer["recommended_actions"]:
            items.append(f"  {action['priority']}. {action['action']} (Owner: {action['owner']})")
    add_bullets(slide, items, size=14)


def save_presentation(prs, path=None):
    if path is not None:
        prs.save(path)
        return path
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def build_presentation_bytes(answer, title=None):
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is required to build PPT reports. Install it with 'pip install python-pptx'.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    global prs_width
    prs_width = prs.slide_width

    subtitle = title or answer.get("question", "EDIP Executive Summary")
    add_title_slide(prs, "EDIP Executive Decision Briefing", subtitle)
    add_answer_summary_slide(prs, answer)

    if answer.get("recommended_actions"):
        slide = add_body_slide(prs, "Recommended Actions")
        lines = [f"{a['priority']}. {a['action']}" for a in answer["recommended_actions"]]
        add_bullets(slide, lines, size=15)

    if answer.get("supporting_evidence"):
        slide = add_body_slide(prs, "Supporting Evidence")
        evidence_lines = [f"{e['source']}: {e['excerpt']}" for e in answer["supporting_evidence"][:5]]
        add_bullets(slide, evidence_lines, size=14)

    return save_presentation(prs)


def build_report():
    consultant = AIConsultant()
    revenue = consultant.answer_revenue_root_cause("Why did revenue drop in Europe in Q2 2026?")
    churn = consultant.answer_churn("Which customers are likely to churn?")
    pricing = consultant.answer_pricing("What pricing strategy should we use in Europe?")
    actions = consultant.answer_next_actions("What should we do next?")

    global prs_width
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs_width = prs.slide_width

    add_title_slide(
        prs,
        "Executive Decision Intelligence Briefing",
        "Q2 2026 Business Review \u2014 Generated by EDIP AI Consultant  |  " + revenue["quantitative_findings"]["latest_month"],
    )

    # --- Slide 2: Revenue root cause ---
    s = add_body_slide(prs, "Why Did Revenue Drop in Q2 2026?")
    qf = revenue["quantitative_findings"]
    add_kpi_cards(s, [
        ("Europe Revenue (latest mo.)", f"${qf['revenue']:,.0f}"),
        ("Month-over-Month Change", f"{qf['month_over_month_change_pct']:+.1f}%"),
        ("Competitor Price Index", f"{qf['competitor_price_index']:.1f}"),
        ("Model Confidence", revenue["confidence"]),
    ], top=Inches(1.3))
    add_bullets(s, [
        f"\u2022 {revenue['narrative']}",
        "",
        "Top model-driven factors (SHAP):",
    ] + [f"    \u2013 {f['feature']}: {f['direction']} predicted revenue (impact {f['shap_contribution']:+.0f})"
         for f in revenue["model_explanation"]["top_contributing_factors"][:3]]
      + ["", "Sources: " + ", ".join(sorted(set(e["source"] for e in revenue["supporting_evidence"])))],
    top=Inches(3.1), size=15)

    # --- Slide 3: Churn ---
    s = add_body_slide(prs, "Customer Churn Risk")
    cs = churn["summary"]
    add_kpi_cards(s, [
        ("Customers Analyzed", f"{cs['customers_analyzed']:,}"),
        ("High Risk (>60%)", f"{cs['high_risk_count']:,}"),
        ("CLV at Risk", f"${cs['clv_at_risk']:,.0f}"),
        ("Avg Churn Prob.", f"{cs['avg_churn_probability']:.1%}"),
    ], top=Inches(1.3))
    add_bullets(s, [churn["narrative"]], top=Inches(3.1), size=15)

    # --- Slide 4: Pricing recommendation ---
    s = add_body_slide(prs, "Pricing Strategy Recommendation")
    add_kpi_cards(s, [
        ("Region", pricing["region_analyzed"]),
        ("Current Avg Price", f"${pricing['current_avg_price']:.2f}"),
        ("Recommended Change", f"{pricing['recommended_price_change_pct']:+.0%}"),
        ("Profit Impact", f"${pricing['projected_profit_impact']:,.0f}/mo"),
    ], top=Inches(1.3))
    add_bullets(s, [pricing["narrative"]], top=Inches(3.1), size=15)

    # --- Slide 5: Recommended actions ---
    s = add_body_slide(prs, "Recommended Actions")
    lines = []
    for a in actions["recommended_actions"]:
        lines.append(f"{a['priority']}. {a['action']}  (Owner: {a['owner']})")
    add_bullets(s, lines, top=Inches(1.3), size=15)

    out_path = OUT_DIR / "EDIP_Executive_Report.pptx"
    prs.save(out_path)
    print(f"Executive report saved to {out_path}")
    return out_path


if __name__ == "__main__":
    build_report()
